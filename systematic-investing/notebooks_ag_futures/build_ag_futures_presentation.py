from __future__ import annotations

import base64
from io import BytesIO
from pathlib import Path

import nbformat
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
NOTEBOOK_PATH = ROOT / "ag_futures_final_project.ipynb"
OUTPUT_PATH = ROOT / "AG_FUTURES_PROJECT_PRESENTATION.pptx"


GREEN = RGBColor(44, 95, 45)
GREEN_DARK = RGBColor(29, 63, 30)
SAND = RGBColor(245, 241, 232)
GOLD = RGBColor(184, 145, 55)
CHARCOAL = RGBColor(34, 34, 34)
GRAY = RGBColor(95, 95, 95)
WHITE = RGBColor(255, 255, 255)


def load_notebook(path: Path):
    return nbformat.read(path, as_version=4)


def get_cell_images(nb, source_snippet: str) -> list[bytes]:
    for cell in nb.cells:
        source = cell.get("source", "")
        if source_snippet in source:
            images = []
            for output in cell.get("outputs", []):
                data = output.get("data", {})
                if "image/png" in data:
                    images.append(base64.b64decode(data["image/png"]))
            return images
    raise KeyError(f"Could not find notebook cell containing: {source_snippet}")


def add_background(slide, title: str | None = None, subtitle: str | None = None):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = SAND

    top = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.8)
    )
    top.fill.solid()
    top.fill.fore_color.rgb = GREEN
    top.line.color.rgb = GREEN

    bottom = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, Inches(7.15), Inches(13.333), Inches(0.35)
    )
    bottom.fill.solid()
    bottom.fill.fore_color.rgb = GREEN_DARK
    bottom.line.color.rgb = GREEN_DARK

    if title:
        tb = slide.shapes.add_textbox(Inches(0.45), Inches(0.12), Inches(9.5), Inches(0.45))
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.name = "Aptos Display"
        run.font.size = Pt(26)
        run.font.bold = True
        run.font.color.rgb = WHITE

    if subtitle:
        tb = slide.shapes.add_textbox(Inches(0.45), Inches(0.42), Inches(10.5), Inches(0.24))
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = subtitle
        run.font.name = "Aptos"
        run.font.size = Pt(11)
        run.font.color.rgb = RGBColor(223, 235, 223)


def add_title_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = GREEN_DARK

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(0.9), Inches(0.18), Inches(4.95)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = GOLD
    accent.line.color.rgb = GOLD

    tb = slide.shapes.add_textbox(Inches(1.05), Inches(1.0), Inches(10.8), Inches(2.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Agricultural Futures\nFinal Project"
    run.font.name = "Aptos Display"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = WHITE

    p2 = tf.add_paragraph()
    p2.space_before = Pt(12)
    run2 = p2.add_run()
    run2.text = (
        "Systematic corn-futures trading with technical strategy diversification\n"
        "and crop-health regime overlays from NASA vegetation data"
    )
    run2.font.name = "Aptos"
    run2.font.size = Pt(17)
    run2.font.color.rgb = RGBColor(226, 236, 226)

    tb3 = slide.shapes.add_textbox(Inches(1.08), Inches(5.85), Inches(10.5), Inches(0.75))
    p3 = tb3.text_frame.paragraphs[0]
    run3 = p3.add_run()
    run3.text = "Robo Advisors and Systematic Trading"
    run3.font.name = "Aptos"
    run3.font.size = Pt(18)
    run3.font.color.rgb = WHITE

    p4 = tb3.text_frame.add_paragraph()
    run4 = p4.add_run()
    run4.text = "James Lawrence"
    run4.font.name = "Aptos"
    run4.font.size = Pt(15)
    run4.font.color.rgb = RGBColor(223, 235, 223)
    return slide


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str], subtitle: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, title, subtitle)

    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.15), Inches(11.8), Inches(5.7))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.space_after = Pt(8)
        p.font.name = "Aptos"
        p.font.size = Pt(22 if i == 0 else 20)
        p.font.color.rgb = CHARCOAL
        if i == 0:
            p.font.bold = True
    return slide


def add_two_column_slide(
    prs: Presentation,
    title: str,
    left_bullets: list[str],
    image_bytes: bytes,
    subtitle: str | None = None,
    image_left: float = 6.9,
    image_top: float = 1.35,
    image_width: float = 5.8,
):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, title, subtitle)

    left = slide.shapes.add_textbox(Inches(0.65), Inches(1.2), Inches(5.8), Inches(5.7))
    tf = left.text_frame
    tf.word_wrap = True
    for i, line in enumerate(left_bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.space_after = Pt(8)
        p.font.name = "Aptos"
        p.font.size = Pt(18)
        p.font.color.rgb = CHARCOAL

    slide.shapes.add_picture(BytesIO(image_bytes), Inches(image_left), Inches(image_top), width=Inches(image_width))
    return slide


def add_chart_slide(prs: Presentation, title: str, image_bytes: bytes, bullets: list[str], subtitle: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, title, subtitle)

    slide.shapes.add_picture(BytesIO(image_bytes), Inches(0.6), Inches(1.2), width=Inches(7.2))

    right = slide.shapes.add_textbox(Inches(8.0), Inches(1.3), Inches(4.7), Inches(5.5))
    tf = right.text_frame
    tf.word_wrap = True
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.space_after = Pt(8)
        p.font.name = "Aptos"
        p.font.size = Pt(18)
        p.font.color.rgb = CHARCOAL
        if i == 0:
            p.font.bold = True
    return slide


def add_table_slide(
    prs: Presentation,
    title: str,
    columns: list[str],
    rows: list[list[str]],
    subtitle: str | None = None,
    footnote: str | None = None,
):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, title, subtitle)

    table = slide.shapes.add_table(
        rows=len(rows) + 1,
        cols=len(columns),
        left=Inches(0.55),
        top=Inches(1.35),
        width=Inches(12.2),
        height=Inches(4.4),
    ).table

    widths = [2.9, 1.4, 1.6, 1.5, 1.8, 1.6]
    for i, width in enumerate(widths[: len(columns)]):
        table.columns[i].width = Inches(width)

    for col_idx, label in enumerate(columns):
        cell = table.cell(0, col_idx)
        cell.text = label
        cell.fill.solid()
        cell.fill.fore_color.rgb = GREEN
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.bold = True
                run.font.color.rgb = WHITE
                run.font.name = "Aptos"
                run.font.size = Pt(14)

    for row_idx, row in enumerate(rows, start=1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if row_idx % 2 else RGBColor(240, 244, 237)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.name = "Aptos"
                    run.font.size = Pt(13)
                    run.font.color.rgb = CHARCOAL

    if footnote:
        tb = slide.shapes.add_textbox(Inches(0.75), Inches(6.0), Inches(11.9), Inches(0.6))
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = footnote
        run.font.name = "Aptos"
        run.font.size = Pt(14)
        run.font.color.rgb = GRAY
    return slide


def build_presentation():
    nb = load_notebook(NOTEBOOK_PATH)
    seasonal_chart = get_cell_images(nb, "fig, axes = plt.subplots(3, 1, figsize=(14, 12), sharex=True)")[0]
    corr_chart = get_cell_images(nb, "corr = raw_optimization_frame.corr()")[0]
    frontier_chart = get_cell_images(nb, "raw_subset_results = optimize_strategy_subsets(raw_optimization_frame)")[0]
    overlay_cum_chart = get_cell_images(nb, "combo_variant_series = {")[0]
    overlay_bar_chart = get_cell_images(nb, "combo_variant_series = {")[1]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)

    add_bullet_slide(
        prs,
        "Project Question",
        [
            "Goal: maximize risk-adjusted performance trading agricultural futures.",
            "Primary market: corn futures, with wheat futures used for a relative-value pairs sleeve.",
            "Core question: can crop-health information improve a diversified systematic strategy portfolio?",
            "Method: build raw technical sleeves first, then test whether realistic weather-forecast skill adds value.",
        ],
        "Motivation and hypothesis",
    )

    add_table_slide(
        prs,
        "Data And Workflow",
        ["Dataset", "Rows", "Start", "End"],
        [
            ["Corn futures", "6,452", "2000-07-17", "2026-04-30"],
            ["Wheat futures", "6,464", "2000-07-17", "2026-04-30"],
            ["NDVI daily panel", "9,546", "2000-02-18", "2026-04-07"],
            ["Corn + NDVI overlap", "6,435", "2000-07-17", "2026-04-07"],
            ["Optimization sample", "6,159", "2001-10-02", "2026-04-07"],
        ],
        subtitle="Prepared caches used by the final notebook",
        footnote="The optimization sample is the common NDVI-overlap window where NDVI, futures data, and all raw strategy return series are simultaneously available.",
    )

    add_chart_slide(
        prs,
        "NASA Signal Engineering",
        seasonal_chart,
        [
            "Raw NDVI level is too seasonal to trade directly.",
            "Each 16-day MODIS composite is reduced to a regional mean and forward-filled to a daily index for alignment.",
            "The final signal is a seasonally adjusted anomaly: current NDVI relative to its historical norm for the same point in the crop calendar.",
            "Only April through October are treated as active crop-health months.",
        ],
        "Why the project moved from raw NDVI level to seasonal anomaly",
    )

    add_bullet_slide(
        prs,
        "Strategy Universe",
        [
            "Trend sleeves: moving-average 10/30, 30/100, and 80/160.",
            "Counter-trend sleeve: buy-the-dip rule based on prior highs and average range.",
            "Volatility sleeve: switch exposure according to the current volatility regime in corn futures.",
            "Pairs sleeves: 5-day, 10-day, and 20-day corn-wheat relative-value trades.",
            "Weather overlay: basic 14-day, advanced 45-day, and expert 90-day crop-health lead windows.",
        ],
        "Raw strategies first, weather overlay second",
    )

    add_table_slide(
        prs,
        "Key Standalone Results",
        ["Sleeve", "Sharpe", "Cum. Return", "Vol", "Max DD"],
        [
            ["vol_regime", "0.4671", "9.6234", "0.3112", "-0.5660"],
            ["ma_10_30", "0.3698", "4.0804", "0.3112", "-0.4664"],
            ["ma_80_160", "0.1728", "0.1339", "0.3113", "-0.8763"],
            ["pairs_10d", "0.0849", "0.0737", "0.1226", "-0.4013"],
            ["countertrend_p2_2", "-0.0183", "-0.4390", "0.1999", "-0.7339"],
            ["ma_30_100", "-0.0985", "-0.8563", "0.3113", "-0.9389"],
        ],
        subtitle="Standalone raw strategy sleeves on the common optimization sample",
        footnote="Corn-wheat spread diagnostics were supportive of a pairs test: ADF statistic = -3.9752, p-value = 0.0015.",
    )

    add_chart_slide(
        prs,
        "Why Combine Strategies?",
        corr_chart,
        [
            "This chart is a correlation matrix of daily raw strategy returns on the common evaluation window.",
            "Low or mildly negative correlation means two sleeves can diversify each other.",
            "The pairs sleeve is largely separate from the outright corn sleeves, which makes it useful at the portfolio level.",
            "This is why the project optimizes combinations of strategies instead of picking one winner.",
        ],
        "Diversification evidence from the raw return streams",
    )

    add_two_column_slide(
        prs,
        "Optimized Raw Combination",
        [
            "The optimizer searched all non-empty subsets of the raw strategy universe and selected long-only weights that maximized annualized Sharpe.",
            "Selected sleeves:",
            "vol_regime: 36.71%",
            "ma_10_30: 28.09%",
            "pairs_10d: 18.37%",
            "ma_80_160: 16.83%",
            "Result: Sharpe 0.6343, cumulative return 6.9749, max drawdown -0.3605.",
            "This outperformed the best raw standalone sleeve (vol_regime at 0.4671 Sharpe).",
        ],
        frontier_chart,
        "Efficient-frontier style search over raw strategy subsets",
        image_left=Inches(6.65) / Inches(1),  # keep numeric form readable below
        image_top=1.35,
        image_width=6.1,
    )

    add_two_column_slide(
        prs,
        "Weather Overlay On The Final Portfolio",
        [
            "The raw portfolio weights were held fixed.",
            "Then the three crop-health forecast tiers were applied as regime overlays.",
            "This isolates the value of weather-model skill from the value of raw strategy diversification.",
            "Only the advanced lead materially improved the final portfolio.",
            "combo_raw_optimal: Sharpe 0.6343",
            "basic: 0.6245 | uplift -0.0098",
            "advanced: 0.7829 | uplift +0.1486",
            "expert: 0.6225 | uplift -0.0118",
        ],
        overlay_bar_chart,
        "Sharpe uplift after applying weather-prediction tiers",
        image_left=6.75,
        image_top=1.45,
        image_width=5.9,
    )

    add_chart_slide(
        prs,
        "Final Interpretation",
        overlay_cum_chart,
        [
            "A diversified systematic portfolio was the main source of performance improvement.",
            "NASA vegetation data only became useful after seasonal adjustment and crop-calendar filtering.",
            "The best result came from medium-horizon crop-health foresight, not from very short or very long leads.",
            "Economic reading: sub-seasonal weather-model skill appears more actionable than either near-term weather or very long-range seasonal guessing.",
        ],
        "What the project suggests about real-world crop-futures trading",
    )

    add_bullet_slide(
        prs,
        "Limitations And Next Steps",
        [
            "The NDVI lead tests are oracle-style benchmarks, not directly tradable forecasts.",
            "Transaction costs, slippage, and margin constraints are not modeled.",
            "The project uses NDVI only; soil moisture, temperature, and precipitation could extend the framework.",
            "Future work: replace oracle NDVI leads with explicit weather-model features and test true out-of-sample forecasting.",
        ],
        "Academic interpretation and realistic extensions",
    )

    prs.save(OUTPUT_PATH)
    print(f"Saved presentation to {OUTPUT_PATH}")


if __name__ == "__main__":
    build_presentation()
